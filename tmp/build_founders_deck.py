#!/usr/bin/env python3
"""
Locore Founders 50 募集デッキを、_template.pptx
（Locore 配色済みテンプレ）のスライドマスター上に生成する。

19 ページ構成は /founders/deck/print と同じ。

Run: python3 tmp/build_founders_deck.py
"""

from __future__ import annotations

import os
import sys
from typing import List, Optional, Tuple

from pptx import Presentation
from pptx.util import Pt

# テンプレ (Locore 配色済) → 完成 .pptx
TEMPLATE = "_template.pptx"
OUTPUT = "Founders50_募集要項.pptx"

# レイアウト index (Master 0)
L_COVER = 0          # （基本版） タイトル_ ロゴ入り_A4
L_COVER_NOLOGO = 1   # （基本版） タイトル ロゴ無_A4
L_BLANK = 5          # 白紙
L_SECTION = 6        # 中表紙 (section divider)
L_TITLE_ONLY = 7     # タイトルのみ (for Stat slide)
L_LEFT = 9           # コンテンツ左サイド
L_TWO_SIDE = 10      # コンテンツ両サイド (for Compare)
L_FULL = 12          # コンテンツ全面 (for headlines / lists)
L_BACK_WHITE = 13    # 背表紙_白 (back cover for CTA)


def clear_text(shape):
    """既存の placeholder のテキストを全消去 (paragraphs 維持)"""
    tf = shape.text_frame
    # 全段落を 1 つに統合 (最初の段落以外を XML 直接削除)
    p0 = tf.paragraphs[0]
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    # 段落 0 の run を全消去
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)


def set_text(shape, text: str, *, size_pt: float | None = None,
             bold: bool | None = None, lines: list[str] | None = None):
    """placeholder にテキストをセット (フォント・色はテーマに従う)"""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.word_wrap = True
    if lines is None:
        lines = text.split('\n')

    # 既存段落をクリアして初期化
    clear_text(shape)
    p0 = tf.paragraphs[0]

    # 各行を段落として追加
    for i, line in enumerate(lines):
        p = p0 if i == 0 else tf.add_paragraph()
        # 段落の既存 run を全消去
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
        run = p.add_run()
        run.text = line
        if size_pt:
            run.font.size = Pt(size_pt)
        if bold is not None:
            run.font.bold = bold


def get_placeholder(slide, idx):
    """placeholder を idx で取得"""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def delete_all_slides(prs):
    """既存の全スライドを完全に削除（XML パッケージから物理削除）

    sldIdLst からの削除だけだと、ppt/slides/slide*.xml が残ったまま
    新規追加と衝突するので、part 自体も drop する。
    """
    # 1. sldIdLst から ID を 1 つずつ取って、対応する slide part を取得 → drop
    slide_id_list = prs.slides._sldIdLst
    # rels: slide_id → rId → slide_part
    rels_to_drop = []
    for sld_id in list(slide_id_list):
        rId = sld_id.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        rels_to_drop.append(rId)
        slide_id_list.remove(sld_id)

    # 2. presentation part の relationship からも slide rel を drop
    pres_part = prs.part
    for rId in rels_to_drop:
        try:
            pres_part.drop_rel(rId)
        except Exception:
            pass


def add_slide(prs, layout_idx):
    layout = prs.slide_masters[0].slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


# ============================================================================
# スライド生成ヘルパ
# ============================================================================

def slide_cover(prs, title: str, subtitle: str, kicker: str = ""):
    """L_COVER_NOLOGO を使った表紙 (ロゴは現状未配置なので no-logo 版を選択)"""
    s = add_slide(prs, L_COVER_NOLOGO)
    if (t := get_placeholder(s, 0)):
        set_text(t, title, size_pt=44, bold=True)
    if (t := get_placeholder(s, 1)):
        set_text(t, subtitle, size_pt=18)
    # date を kicker として使う
    if kicker and (t := get_placeholder(s, 12)):
        set_text(t, kicker, size_pt=11)
    return s


def slide_section(prs, chapter: int, title: str, oneliner: str = ""):
    s = add_slide(prs, L_SECTION)
    big = f"Chapter {chapter:02d}  /  {title}"
    if oneliner:
        big += f"\n{oneliner}"
    if (t := get_placeholder(s, 10)):
        # 章タイトル + 一行説明を一つの placeholder に
        set_text(t, big, size_pt=32, bold=True,
                 lines=[f"Chapter {chapter:02d}", title, "", oneliner])
    return s


def slide_headline(prs, header: str, key_message: str, body: list[str]):
    """L_FULL を使った標準コンテンツ"""
    s = add_slide(prs, L_FULL)
    # Header (small top)
    if (t := get_placeholder(s, 15)):
        set_text(t, header, size_pt=14)
    # Key message (big title)
    if (t := get_placeholder(s, 0)):
        set_text(t, key_message, size_pt=24, bold=True)
    # Body (object)
    if (t := get_placeholder(s, 1)):
        set_text(t, '\n'.join(body), size_pt=14, lines=body)
    return s


def slide_bullet_list(prs, header: str, key_message: str,
                      items: list[tuple[str, str]]):
    """箇条書きリスト: items は (title, description) のタプル"""
    s = add_slide(prs, L_FULL)
    if (t := get_placeholder(s, 15)):
        set_text(t, header, size_pt=14)
    if (t := get_placeholder(s, 0)):
        set_text(t, key_message, size_pt=24, bold=True)
    if (t := get_placeholder(s, 1)):
        lines = []
        for i, (title, desc) in enumerate(items):
            lines.append(f"{i+1:02d}.  {title}")
            if desc:
                lines.append(f"      {desc}")
            lines.append("")
        # 最後の空行を削除
        if lines and lines[-1] == "":
            lines.pop()
        set_text(t, "", size_pt=13, lines=lines)
    return s


def slide_compare(prs, header: str, key_message: str,
                  left_title: str, left_body: list[str],
                  right_title: str, right_body: list[str]):
    s = add_slide(prs, L_TWO_SIDE)
    if (t := get_placeholder(s, 15)):
        set_text(t, header, size_pt=14)
    if (t := get_placeholder(s, 0)):
        set_text(t, key_message, size_pt=24, bold=True)
    # 左カラム
    if (t := get_placeholder(s, 1)):
        lines = [left_title, ""] + [f"・ {b}" for b in left_body]
        set_text(t, "", size_pt=13, lines=lines)
    # 右カラム
    if (t := get_placeholder(s, 12)):
        lines = [right_title, ""] + [f"・ {b}" for b in right_body]
        set_text(t, "", size_pt=13, lines=lines)
    return s


def slide_stat(prs, header: str, value: str, label: str, caption: str = ""):
    """L_TITLE_ONLY で巨大な数値を中央に置く"""
    s = add_slide(prs, L_TITLE_ONLY)
    if (t := get_placeholder(s, 15)):
        set_text(t, header, size_pt=14)
    if (t := get_placeholder(s, 0)):
        # value + label を 1 つの placeholder に
        set_text(t, "", size_pt=72, bold=True,
                 lines=[value, "", label, "", caption])
    return s


def slide_cta(prs, title: str, cta_text: str, url: str, hint: str = ""):
    """L_BACK_WHITE をクロージング CTA に使う"""
    s = add_slide(prs, L_BACK_WHITE)
    if (t := get_placeholder(s, 10)):
        set_text(t, "", size_pt=28, bold=True,
                 lines=[title, "", cta_text, url, "", hint])
    return s


# ============================================================================
# Main: 19 ページ生成
# ============================================================================

def main():
    if not os.path.exists(TEMPLATE):
        print(f"Template not found: {TEMPLATE}", file=sys.stderr)
        sys.exit(1)

    # テンプレを開いて既存スライドを全消去 (レイアウトは温存)
    prs = Presentation(TEMPLATE)
    delete_all_slides(prs)

    # ─── 1. Cover ─────────────────────────────────────────────
    slide_cover(
        prs,
        title="現地民が語る、もう一段深い旅。",
        subtitle="駐在員が書く、信頼できる旅行誌。\nLocore は最初の書き手 50 人を募集します。",
        kicker="Founders 50 募集要項  /  Locore",
    )

    # ─── 2. Mission Quote ─────────────────────────────────────
    slide_section(
        prs,
        chapter=0,
        title="ローカル × 語り部 × コア。",
        oneliner="その街の輪郭を、住人の言葉で。",
    )

    # ─── 3. Section 01 ───────────────────────────────────────
    slide_section(
        prs,
        chapter=1,
        title="なぜ Locore を作るのか",
        oneliner="既存の旅行情報の課題と、駐在員の眠っている知見について。",
    )

    # ─── 4. Problem 1 ─────────────────────────────────────────
    slide_headline(
        prs,
        header="課題 ①",
        key_message="旅行情報は、もう同じ顔ばかり。",
        body=[
            "ガイドブックも、SNS も、検索結果も、似たような観光地の似たような写真ばかりが並ぶ。",
            "AI で 5 分で出てくる情報を、わざわざお金を払って読む人はいない。",
            "",
            "既存メディアは広告依存で匿名性が高く、誰がどんな立場で書いた情報なのかが見えづらい。",
            "結果、混雑する場所に観光客が集中し、本当に良い場所は埋もれていく。",
        ],
    )

    # ─── 5. Problem 2 ─────────────────────────────────────────
    slide_headline(
        prs,
        header="課題 ②",
        key_message="駐在員の知見が、SNS に流れて消えていく。",
        body=[
            "海外に住む日本人は、現地にしかない知識を毎日蓄積している。",
            "でも Instagram のストーリーは 24 時間で消え、ブログは技術的負担が大きく、",
            "YouTube は収益の不確実性が高い。",
            "",
            "結果、価値ある情報が「個人の善意」だけに支えられて流れていく。",
            "発信者には対価が回らず、読み手は信頼できる情報源にたどり着けない。",
        ],
    )

    # ─── 6. Section 02 ───────────────────────────────────────
    slide_section(
        prs,
        chapter=2,
        title="Locore の解決策",
        oneliner="一次情報・モデルコース・ローカルスコアの 3 つで、画一化を超える。",
    )

    # ─── 7. 3 つの設計原則 ──────────────────────────────────
    slide_bullet_list(
        prs,
        header="3 つの設計原則",
        key_message="Locore は、これだけはやらない。",
        items=[
            ("一次情報・自撮しか載せない",
             "投稿者本人が現地で撮った写真、現地で体験した内容のみ。AI 生成画像と、検索で出てくる一般情報は掲載対象外。"),
            ("スポット単発ではなくモデルコース",
             "時間帯・混雑・移動難易度まで含めた「実際に辿れるルート」を提示。旅行者の調査負担と不安を減らす。"),
            ("ローカルスコアで選別",
             "記事はローカル度を採点。観光地の単純紹介は低評価・修正対象。住人だけが知る場所が上位に来る設計。"),
        ],
    )

    # ─── 8. ローカルスコア ────────────────────────────────────
    slide_headline(
        prs,
        header="ローカルスコア",
        key_message="ローカル性が高いほど、上位に。",
        body=[
            "すべての記事は「どれだけローカルか」のスコアで評価される。",
            "Google マップでレビュー数が多い観光地は低評価、住人だけが知る場所や時間帯は高評価。",
            "",
            "観光地の「裏技」(混まない時間帯、地元の楽しみ方) は許容するが、単なる紹介は修正対象。",
            "理念に共感した書き手だけが残る仕組み。",
        ],
    )

    # ─── 9. Compare ──────────────────────────────────────────
    slide_compare(
        prs,
        header="何が違うのか",
        key_message="既存メディアとの違い",
        left_title="既存の旅行メディア（広告依存・匿名・一律）",
        left_body=[
            "匿名アカウント中心で誰が書いたか不明",
            "PR と本物が区別できない",
            "\"映え\" 優先で実用性が薄い",
            "広告がコンテンツに混入",
        ],
        right_title="Locore（居住確認済みの書き手だけ）",
        right_body=[
            "対象都市に 1 年以上住む書き手のみ",
            "編集チームが事実確認",
            "広告なし。読者課金で運営",
            "一次情報・自撮・モデルコース重視",
        ],
    )

    # ─── 10. 3 つの読者層 ───────────────────────────────────
    slide_bullet_list(
        prs,
        header="Market",
        key_message="3 つの読者層に届ける。",
        items=[
            ("観光地を避けたい旅行者",
             "混雑する有名観光地ではなく、ローカルが普通に通う場所を求める層。年に 1〜2 回の海外旅行で、調査に時間をかけたくない人。"),
            ("海外を志向する若手",
             "留学・駐在・ワーホリを検討中の人。SNS では聞きづらい本音の話を、近い境遇の書き手から直接聞きたい層。"),
            ("現地在住の駐在員・留学生",
             "日々の発信を価値化したい人。技術的負担なく書け、執筆と提供サービスの両方で収益機会を得たい層。"),
        ],
    )

    # ─── 11. Section 03 ──────────────────────────────────────
    slide_section(
        prs,
        chapter=3,
        title="ビジネスモデル",
        oneliner="広告に依存せず、読者と書き手の信頼で成り立つ仕組み。",
    )

    # ─── 12. 3 つの収益源 ───────────────────────────────────
    slide_bullet_list(
        prs,
        header="3 つの収益源",
        key_message="広告は載せない。",
        items=[
            ("記事の決済手数料",
             "1 本 ¥600〜¥1,500 の有料記事。決済手数料の一部を運営費・編集チームの人件費に充当。"),
            ("サブスクリプション",
             "月額会員向けの限定記事・先行公開・特集アクセス。コア読者からの安定収益を作る。"),
            ("ユーザー提供サービス（無手数料）",
             "駐在員が個別に提供する相談・ガイド・撮影アテンドなど。プラットフォーム手数料は取らず、書き手のビジネス機会を最大化。"),
        ],
    )

    # ─── 13. Section 04 ──────────────────────────────────────
    slide_section(
        prs,
        chapter=4,
        title="Founders 50",
        oneliner="最初の 50 人を、長く覚えていたい。",
    )

    # ─── 14. Stat 50 ─────────────────────────────────────────
    slide_stat(
        prs,
        header="Founders",
        value="50 人",
        label="先着 50 人、限定で。",
        caption="認証バッジ、優遇手数料、方針への発言権。すべて帰任後も続く恒久的な特典。",
    )

    # ─── 15. Founders 特典 ──────────────────────────────────
    slide_bullet_list(
        prs,
        header="参加するとどうなるか",
        key_message="3 つの恒久的特典。",
        items=[
            ("認証バッジ",
             "名前の横に「Founders」表示。読み手から見て最初に信頼される目印。プロフィールにも常時表示。"),
            ("優遇手数料",
             "ライターの取り分が通常 75% のところ、Founders は最終的に 90% に。条件は帰任後も継続。"),
            ("サービス方針への発言権",
             "機能追加・コミュニティルール・新都市展開などに、Founders の意見が反映される仕組み。"),
        ],
    )

    # ─── 16. 参加条件 ────────────────────────────────────────
    slide_headline(
        prs,
        header="参加条件",
        key_message="在住 1 年以上、月 5 本、理念への共感。",
        body=[
            "対象都市に 1 年以上居住している方。月 5 本以上の執筆コミットメントをお願いしています。",
            "記事の質はローカルスコアと編集チームの審査を通過する必要があります。",
            "",
            "初期 (ローンチ前) は 1 本 2,000 円で買い取り、ローンチ後は売上連動に移行。",
            "観光地の「裏技」は OK ですが、単純な紹介記事は修正対象です。",
        ],
    )

    # ─── 17. 応募の流れ ─────────────────────────────────────
    slide_bullet_list(
        prs,
        header="Application",
        key_message="応募の流れ",
        items=[
            ("応募フォームの提出",
             "自己紹介、希望都市・テーマ、過去の記事サンプル、SNS リンク、居住証明（公共料金請求書や賃貸契約書）を提出。"),
            ("書類選考",
             "理念への共感、現地知見、執筆の継続力を判断。3〜5 営業日でご連絡。"),
            ("Zoom 面談（15〜30 分）",
             "編集チームと簡単な対話。書きたいテーマ、Locore に期待することを伺います。"),
            ("結果通知 → 執筆開始",
             "採用された方には Founders オンボーディングを案内。"),
        ],
    )

    # ─── 18. ロードマップ + チーム ────────────────────────
    slide_headline(
        prs,
        header="Roadmap & Team",
        key_message="2026 年 9 月公開、フランスから。",
        body=[
            "公開目標は 2026 年 9 月 30 日。Founders 募集は 50 名に到達次第締切。",
            "初期展開はフランス（パリ・ボルドー等）、その後ヨーロッパ各都市、将来的に日本展開も視野に。",
            "",
            "運営チームは現在、創業者（チェコ在住・フランス留学経験・コンサル業）と",
            "エンジニアの 2 名で進行中。各拠点に編集・運営ヘッドを順次採用予定。",
        ],
    )

    # ─── 19. CTA ─────────────────────────────────────────────
    slide_cta(
        prs,
        title="あなたの街を、書きませんか。",
        cta_text="応募する",
        url="locore.app/founders",
        hint="先着 50 人で締切。応募 → 書類選考 → Zoom 面談 → 結果通知。\n"
             "ご質問は hello@locore.app まで。",
    )

    # 保存
    if os.path.exists(OUTPUT):
        os.remove(OUTPUT)
    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")
    print(f"Size: {os.path.getsize(OUTPUT)} bytes")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
